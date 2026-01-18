from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define some colors
    GREEN = RGBColor(21, 128, 61)  # #15803d
    BLUE = RGBColor(14, 165, 233)   # #0ea5e9
    DARK_TEXT = RGBColor(30, 41, 59)

    def add_slide(title_text, content_items=None, layout_idx=1):
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        
        # Title styling
        title = slide.shapes.title
        title.text = title_text
        title_para = title.text_frame.paragraphs[0]
        title_para.font.bold = True
        title_para.font.color.rgb = GREEN
        title_para.font.size = Pt(36)

        if content_items:
            tf = slide.placeholders[1].text_frame
            tf.word_wrap = True
            for item in content_items:
                p = tf.add_paragraph()
                p.text = item
                p.level = 0
                p.space_after = Pt(10)
                p.font.size = Pt(18)
                p.font.color.rgb = DARK_TEXT

        return slide

    # --- Slide 1: Team Details ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "AgriBid: TechSprint Project"
    subtitle.text = "Team Details:\na. Team Name: AgriBid\nb. Team Leader Name: [Enter Name]\nc. Problem Statement: Open Innovation"
    
    title_para = title.text_frame.paragraphs[0]
    title_para.font.color.rgb = GREEN
    title_para.font.bold = True

    # --- Slide 2: Brief about solution ---
    add_slide(
        "Brief about your solution and problem statement addressing",
        [
            "Problem: Farmers face a 30-40% profit loss due to non-transparent pricing and middlemen intermediaries in the agricultural supply chain.",
            "Solution: AgriBid is a direct decentralized marketplace that enables competitive bidding between retailers and farmers.",
            "Impact: Eliminates middlemen, ensures real-time price discovery, and provides AI-driven market intelligence to maximize farmer income."
        ]
    )

    # --- Slide 3: Opportunities ---
    add_slide(
        "Opportunities",
        [
            "a. How different is it? Unlike static e-commerce, we offer real-time Mandi-synced AI predictions and a voice-assisted interface for non-digital natives.",
            "b. How will it solve the problem? By creating a direct, transparent bidding war for crops, farmers get the true market value while retailers get fresher produce at lower overheads."
        ]
    )

    # --- Slide 4: List of features ---
    add_slide(
        "List of features offered by the solution",
        [
            "• AI Market Intelligence Hub (Price Predictions)",
            "• Multi-lingual Voice Assistance (Hindi/Kannada)",
            "• Dynamic Bidding (Mini & Bulk Bidding Cycles)",
            "• Performance Dashboard (Sales & Bidding Analytics)",
            "• Secure Digital Wallet for instant settlements"
        ]
    )

    # --- Slide 5: Google Technologies ---
    add_slide(
        "Google Technologies used in the solution",
        [
            "• Gemini 3 Flash: Core AI for processing Mandi data and generating actionable insights.",
            "• Google Cloud Platform: Hosting the marketplace infrastructure for scale and security.",
            "• Google Text-to-Speech: Powering our accessibility layer for farmer-friendly interaction.",
            "• Firebase: Used for real-time notifications and bidding updates."
        ]
    )

    # --- Slide 6: Process flow diagram ---
    add_slide(
        "Process flow diagram or Use-case diagram",
        [
            "1. Farmer Registration & Verification",
            "2. Crop Listing with Quantity/Base Price",
            "3. AI Engine provides Price Suggestions",
            "4. Retailers browse and place Competitive Bids",
            "5. Bid Completion & Automated Payment via Wallet",
            "6. Direct Fulfillment/Delivery orchestration"
        ]
    )

    # --- Slide 7: Wireframes/Mock diagrams ---
    add_slide(
        "Wireframes/Mock diagrams of the proposed solution (optional)",
        [
            "• Landing Page: Clean, high-contrast UI for outdoor usage visibility.",
            "• Bidding Panel: Real-time countdown timers and bid increment controls.",
            "• Market Hub: Interactive Recharts showing price volatility and AI forecasts."
        ]
    )

    # --- Slide 8: Architecture diagram ---
    add_slide(
        "Architecture diagram of the proposed solution",
        [
            "• Frontend: Next.js (React 19) + Lucide Icons + Tailwind Styling",
            "• Backend: Next.js API Routes (Node.js) + Prisma ORM",
            "• Database: PostgreSQL (Neon / GCP Cloud SQL)",
            "• Inference: Google Vertex AI / Gemini 3 Flash API Integration"
        ]
    )

    # --- Slide 9: Snapshots of the MVP ---
    add_slide(
        "Snapshots of the MVP",
        [
            "[Feature 1]: Live AI Market Insights Page with Predicted Prices.",
            "[Feature 2]: Farmer Dashboard showing active bids in Kannada/Hindi.",
            "[Feature 3]: Retailer Marketplace with verified crop listings.",
            "[Feature 4]: Secure Wallet transaction history screen."
        ]
    )

    # --- Slide 10: Additional Details ---
    add_slide(
        "Additional Details/Future Development (if any)",
        [
            "• Logistics Integration: Smart routing for farm-to-retailer transportation.",
            "• Computer Vision: Using Gemini 3 Flash for quality grading via crop photos.",
            "• Financial Inclusion: Micro-credit facilities based on bidding history."
        ]
    )

    # --- Slide 11: Provide links ---
    add_slide(
        "Provide links to your:",
        [
            "1. GitHub Public Repository: [Insert Link]",
            "2. Demo Video Link (3 Minutes): [Insert Link]",
            "3. MVP Link: [Insert Link]"
        ]
    )

    # Save the presentation
    file_path = "AgriBid_TechSprint_Presentation.pptx"
    prs.save(file_path)
    print(f"Presentation saved to {file_path}")

if __name__ == "__main__":
    create_presentation()
